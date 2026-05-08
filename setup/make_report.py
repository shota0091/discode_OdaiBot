"""AI活用開発事例報告書 PowerPoint 生成スクリプト"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_PRIMARY = RGBColor(0x58, 0x65, 0xF2)
C_DARK    = RGBColor(0x2F, 0x31, 0x36)
C_ACCENT  = RGBColor(0x3B, 0xA5, 0x5C)
C_DANGER  = RGBColor(0xED, 0x42, 0x45)
C_WARN    = RGBColor(0xFA, 0xA6, 0x1A)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xF2, 0xF3, 0xF5)
C_TEXT    = RGBColor(0x2E, 0x33, 0x38)
C_MUTED   = RGBColor(0x72, 0x76, 0x7D)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tx(slide, l, t, w, h, text, size, bold=False, color=C_TEXT,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box

def header(slide, title, color=C_PRIMARY):
    rect(slide, 0, 0, 13.33, 1.1, color)
    tx(slide, 0.5, 0.2, 12, 0.7, title, 26, bold=True, color=C_WHITE)

def footer(slide, text="社内勉強会資料 / AI活用開発事例 2026.04"):
    rect(slide, 0, 7.1, 13.33, 0.4, C_DARK)
    tx(slide, 0.3, 7.12, 12, 0.3, text, 9, color=C_MUTED)

def card(slide, l, t, w, h, title, title_color=C_PRIMARY, bg_color=C_WHITE):
    rect(slide, l, t, w, h, bg_color)
    rect(slide, l, t, w, 0.45, title_color)
    tx(slide, l+0.1, t+0.05, w-0.2, 0.38, title, 12, bold=True, color=C_WHITE)

# ════════════════════════════════════════════
# Slide 1: タイトル
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_DARK)
rect(sl, 0, 2.9, 13.33, 0.07, C_PRIMARY)
tx(sl, 1, 0.8, 11, 1.0, "Discord お題Bot 改修プロジェクト", 34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tx(sl, 1, 1.9, 11, 0.6, "AI（Claude）を活用した個人開発事例報告", 20, color=RGBColor(0xB9,0xBB,0xBE), align=PP_ALIGN.CENTER)
tx(sl, 1, 2.6, 11, 0.45, "2026年4月　堀内", 14, color=C_MUTED, align=PP_ALIGN.CENTER)
rect(sl, 3.8, 3.3, 5.7, 0.75, C_PRIMARY)
tx(sl, 3.8, 3.3, 5.7, 0.75, "要件定義・レビュー担当：堀内（SE）", 13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tx(sl, 1, 6.5, 11, 0.4, "※ 個人の趣味開発における AI 活用事例です", 10, color=C_MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# Slide 2: アジェンダ
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "アジェンダ")
items = [
    ("01", "お題Botとは？（システム概要）"),
    ("02", "旧システムの課題"),
    ("03", "改修要件（7項目）"),
    ("04", "新システム構成・Dashboard操作フロー"),
    ("05", "AI活用の役割分担と注意点"),
    ("06", "素人がAIだけで作る危険性"),
    ("07", "工数削減効果"),
    ("08", "まとめ・今後の展望"),
]
for i, (num, label) in enumerate(items):
    col = i % 2
    row = i // 2
    lx = 0.5 + col * 6.4
    ty = 1.4 + row * 1.3
    rect(sl, lx, ty, 6.0, 1.05, C_WHITE)
    rect(sl, lx, ty, 0.65, 1.05, C_PRIMARY)
    tx(sl, lx, ty, 0.65, 1.05, num, 16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tx(sl, lx+0.75, ty+0.25, 5.1, 0.55, label, 13, bold=True, color=C_TEXT)
footer(sl)


# ════════════════════════════════════════════
# Slide 3: お題Botとは？
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "01  お題Botとは？（システム概要）")

# 概要説明
rect(sl, 0.4, 1.25, 12.5, 1.1, C_WHITE)
tx(sl, 0.6, 1.3, 12, 0.45,
   "Discord サーバー向けのお題画像自動投稿 Bot。管理者がDashboardでお題を登録・設定し、Botが定期的にチャンネルへ投稿する。",
   12, color=C_TEXT)
tx(sl, 0.6, 1.72, 12, 0.45,
   "チャンネルごとに独立したローテーション管理を行い、全お題を使い切ったら自動リセットして再投稿する仕組み。",
   12, color=C_TEXT)

# 3機能カード
features = [
    ("お題投稿", "C_PRIMARY", [
        "/odai コマンドで即時投稿",
        "未投稿お題をランダム選択",
        "全使用済みで自動リセット",
        "チャンネル別に独立管理",
    ]),
    ("定期通知", "C_ACCENT", [
        "スケジュール設定で自動投稿",
        "時刻（HH:MM）で毎日実行",
        "タグでお題を絞り込み可能",
        "複数チャンネル・複数時刻対応",
    ]),
    ("Dashboard管理", "C_DARK", [
        "お題画像の登録・削除・編集",
        "タグ管理でカテゴリ分け",
        "スケジュール設定",
        "ユーザー管理（招待制）",
    ]),
]
colors_map = {"C_PRIMARY": C_PRIMARY, "C_ACCENT": C_ACCENT, "C_DARK": C_DARK}
for j, (title, ckey, lines) in enumerate(features):
    lx = 0.4 + j * 4.3
    col = colors_map[ckey]
    card(sl, lx, 2.55, 4.0, 4.2, title, title_color=col)
    for k, line in enumerate(lines):
        tx(sl, lx+0.15, 3.1 + k*0.68, 3.7, 0.6, f"• {line}", 11, color=C_TEXT)

footer(sl)


# ════════════════════════════════════════════
# Slide 4: システム全体像
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "01  システム全体像")

# コンポーネント図（テキストベース）
components = [
    (0.4,  2.0, 2.5, 1.2, "Discord\nサーバー",        RGBColor(0x57,0x86,0xC1)),
    (3.5,  2.0, 2.5, 1.2, "Discord Bot\n(OdaiBot)",   C_DARK),
    (6.6,  1.2, 2.5, 1.2, "FastAPI\n(REST API)",      C_PRIMARY),
    (6.6,  2.9, 2.5, 1.2, "Dashboard\n(Web UI)",      RGBColor(0x8B,0x5C,0xF6)),
    (9.7,  2.0, 2.5, 1.2, "MySQL\nデータベース",      C_ACCENT),
    (3.5,  4.5, 2.5, 1.0, "管理者\n(Discord)",        C_MUTED),
    (6.6,  4.5, 2.5, 1.0, "管理者\n(Dashboard)",      C_MUTED),
]
for lx, ty, w, h, label, color in components:
    rect(sl, lx, ty, w, h, color)
    tx(sl, lx, ty, w, h, label, 12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 矢印テキスト
arrows = [
    (2.95, 2.5,  "→ 投稿"),
    (6.05, 1.7,  "→ API"),
    (9.15, 1.7,  "→ DB"),
    (9.15, 3.4,  "→ DB"),
    (6.05, 3.4,  "← 取得"),
    (3.5,  4.1,  "/odaiコマンド"),
    (6.6,  4.1,  "HTTPS アクセス"),
]
for lx, ty, label in arrows:
    tx(sl, lx, ty, 1.6, 0.4, label, 9, color=C_MUTED, align=PP_ALIGN.CENTER)

# 技術スタック
rect(sl, 0.4, 5.8, 12.5, 1.1, C_WHITE)
tx(sl, 0.5, 5.83, 12, 0.35, "技術スタック", 11, bold=True, color=C_PRIMARY)
stacks = [
    ("Bot", "Python 3.12 / discord.py 2.5"),
    ("API", "FastAPI / mysql-connector-python"),
    ("DB",  "MySQL 8.0 / Rocky Linux 9.6 VPS"),
    ("UI",  "Vanilla JS SPA / Cloudflare Tunnel(HTTPS)"),
]
for j, (tag, val) in enumerate(stacks):
    lx = 0.5 + j * 3.1
    rect(sl, lx, 6.2, 0.6, 0.4, C_PRIMARY)
    tx(sl, lx, 6.2, 0.6, 0.4, tag, 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tx(sl, lx+0.65, 6.22, 2.4, 0.35, val, 9, color=C_TEXT)
footer(sl)


# ════════════════════════════════════════════
# Slide 5: 旧システムの課題
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "02  旧システムの課題", C_DANGER)

rect(sl, 0.4, 1.25, 12.5, 1.7, C_WHITE)
tx(sl, 0.5, 1.3, 12, 0.4, "旧構成", 12, bold=True, color=C_DANGER)
old_boxes = [
    (0.5,  1.75, 2.0, 0.85, "Discord Bot\n(単体)", C_DARK),
    (3.1,  1.75, 2.2, 0.85, "JSON\n({id}_odai.json)", C_DANGER),
    (5.9,  1.75, 2.2, 0.85, "画像ファイル\n(templates/{id}/)", C_DANGER),
    (8.7,  1.75, 2.5, 0.85, "Schedule JSON\n({id}_Schedule.json)", C_DANGER),
]
for lx, ty, w, h, label, color in old_boxes:
    rect(sl, lx, ty, w, h, color)
    tx(sl, lx, ty, w, h, label, 10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
for x in [2.55, 5.35, 8.15]:
    tx(sl, x, 2.1, 0.5, 0.4, "->", 14, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

problems = [
    ("管理の煩雑さ",     "画像とJSONが別管理。整合性が崩れるリスクが常にあった"),
    ("ローテーション",   "投稿フラグがグローバル管理。チャンネルごとの独立管理が不可能"),
    ("管理画面なし",     "設定変更・お題追加にコード修正とサーバー操作が毎回必要"),
    ("モバイル未対応",   "PCのみ想定の設計。スマホでの操作・確認が困難"),
]
for j, (title, desc) in enumerate(problems):
    col = j % 2
    row = j // 2
    lx = 0.4 + col*6.4
    ty = 3.2 + row*1.5
    rect(sl, lx, ty, 6.1, 1.35, C_WHITE)
    rect(sl, lx, ty, 0.18, 1.35, C_DANGER)
    tx(sl, lx+0.28, ty+0.1,  5.7, 0.45, title, 13, bold=True, color=C_DANGER)
    tx(sl, lx+0.28, ty+0.62, 5.7, 0.55, desc,  11, color=C_TEXT)
footer(sl)


# ════════════════════════════════════════════
# Slide 6: 改修要件7項目
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "03  改修要件（7項目）")
reqs = [
    ("①", "プレビュー表示修正",  "お題画像がモーダルからはみ出す問題を解消。CSS制御で枠内に収める"),
    ("②", "ID列の非表示化",      "テーブルのID列を完全削除。内部処理にのみ使用しUIをシンプル化"),
    ("③", "ファイル名あいまい検索","ファイル名のキーワード検索（SQL LIKE）を実装。探しやすさを改善"),
    ("④", "スマホ対応",           "モバイルでのモーダルサイズ・レイアウトを修正。ボトムシート形式に"),
    ("⑤", "横スクロール廃止",    "お題・スケジュール管理の横スクロールを解消。非本質列を非表示に"),
    ("⑥", "自動リセット機能",    "全お題使用済み時にodai_usageを自動削除してリセット・再投稿"),
    ("⑦", "チャンネル別ローテーション","odai_usageテーブル導入。チャンネルAのリセットがBに影響しない独立管理"),
]
for j, (num, title, desc) in enumerate(reqs):
    col = j % 2 if j < 6 else 0
    row = j // 2
    lx = 0.4 + col*6.45
    ty = 1.3 + row*1.28
    w = 12.5 if j == 6 else 6.1
    rect(sl, lx, ty, w, 1.15, C_WHITE)
    rect(sl, lx, ty, 0.55, 1.15, C_PRIMARY)
    tx(sl, lx, ty, 0.55, 1.15, num, 16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tx(sl, lx+0.65, ty+0.08, w-0.75, 0.42, title, 12, bold=True, color=C_PRIMARY)
    tx(sl, lx+0.65, ty+0.56, w-0.75, 0.52, desc,  10, color=C_TEXT)
footer(sl)


# ════════════════════════════════════════════
# Slide 7: Dashboard 操作フロー
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "04  Dashboard 操作フロー（画面イメージ）")

screens = [
    ("ログイン画面", C_DARK, [
        "ユーザー名 / PW 入力",
        "JWT トークン認証",
        "複数サーバー対応",
    ]),
    ("ダッシュボード", C_PRIMARY, [
        "お題数 / 使用済み数",
        "スケジュール数",
        "クイックリンク",
    ]),
    ("お題管理", RGBColor(0x8B,0x5C,0xF6), [
        "画像一覧・プレビュー",
        "ファイル名検索",
        "タグ管理・一括操作",
    ]),
    ("スケジュール管理", C_ACCENT, [
        "チャンネル選択（予測変換）",
        "時刻・タグモード設定",
        "有効/無効切り替え",
    ]),
]

# フロー矢印
for j, (title, color, lines) in enumerate(screens):
    lx = 0.35 + j*3.25
    # 画面モックアップ枠
    rect(sl, lx, 1.25, 2.95, 4.5, C_WHITE)
    rect(sl, lx, 1.25, 2.95, 0.45, color)
    tx(sl, lx, 1.25, 2.95, 0.45, title, 11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # モック内容
    rect(sl, lx+0.1, 1.8, 2.75, 1.6, C_LIGHT)
    tx(sl, lx+0.2, 1.85, 2.55, 1.5,
       "┌─────────┐\n│  " + title + "   │\n│             │\n└─────────┘",
       8, color=C_MUTED)
    for k, line in enumerate(lines):
        rect(sl, lx+0.1, 3.55+k*0.55, 2.75, 0.48, C_LIGHT)
        tx(sl, lx+0.2, 3.6+k*0.55, 2.55, 0.4, f"• {line}", 10, color=C_TEXT)
    if j < 3:
        tx(sl, lx+3.05, 3.1, 0.3, 0.5, ">", 20, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# 初回セットアップフロー
rect(sl, 0.35, 5.95, 12.6, 1.0, RGBColor(0xEE,0xF0,0xFF))
tx(sl, 0.5, 5.98, 12, 0.35, "初回セットアップフロー", 11, bold=True, color=C_PRIMARY)
flow = ["Discord で\n/odai_dashboard", "招待URLを\n受け取る", "Dashboard で\nPW設定・登録", "お題を\nアップロード", "スケジュールを\n設定", "自動投稿\nスタート!"]
for j, step in enumerate(flow):
    lx = 0.5 + j*2.1
    rect(sl, lx, 6.35, 1.9, 0.55, C_PRIMARY if j == 5 else C_DARK)
    tx(sl, lx, 6.35, 1.9, 0.55, step, 9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    if j < 5:
        tx(sl, lx+1.95, 6.55, 0.2, 0.3, ">", 10, bold=True, color=C_PRIMARY)
footer(sl)


# ════════════════════════════════════════════
# Slide 8: Before / After
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "04  新システム構成（Before / After）")

rect(sl, 0.4, 1.25, 5.9, 5.5, C_WHITE)
rect(sl, 0.4, 1.25, 5.9, 0.5, C_DANGER)
tx(sl, 0.4, 1.25, 5.9, 0.5, "  Before（旧構成）", 14, bold=True, color=C_WHITE)
before = [
    "Discord Bot（単体スクリプト）",
    "  ↓ 読み書き",
    "JSONファイル（お題リスト）",
    "  ↓ 参照",
    "画像ファイル（ローカルFS）",
    "  ↓ 参照",
    "スケジュールJSON",
    "",
    "管理方法：コード修正のみ",
    "モバイル：非対応",
    "HTTPS：なし",
]
for j, line in enumerate(before):
    color = C_DANGER if "JSON" in line or "画像" in line else C_TEXT
    tx(sl, 0.6, 1.85+j*0.38, 5.5, 0.36, line, 11, color=color)

tx(sl, 6.25, 3.7, 0.7, 0.7, "=>", 22, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

rect(sl, 6.9, 1.25, 6.0, 5.5, C_WHITE)
rect(sl, 6.9, 1.25, 6.0, 0.5, C_ACCENT)
tx(sl, 6.9, 1.25, 6.0, 0.5, "  After（新構成）", 14, bold=True, color=C_WHITE)
after_items = [
    ("Discord Bot（OdaiBot）",          C_DARK),
    ("FastAPI REST API",                 C_PRIMARY),
    ("MySQL Database",                   C_ACCENT),
    ("  odai / tags / schedules テーブル", C_MUTED),
    ("  odai_usage（チャンネル別管理）",   C_MUTED),
    ("Dashboard（SPA / Vanilla JS）",    RGBColor(0x8B,0x5C,0xF6)),
    ("Cloudflare Tunnel（HTTPS）",       RGBColor(0xF4,0x81,0x20)),
    ("",                                 C_WHITE),
    ("管理方法：Dashboard から操作",      C_ACCENT),
    ("モバイル：対応済み",               C_ACCENT),
    ("HTTPS：Cloudflare Tunnel",         C_ACCENT),
]
for j, (item, color) in enumerate(after_items):
    if color in (C_MUTED, C_TEXT, C_ACCENT) and not item.startswith("管理"):
        tx(sl, 7.1, 1.85+j*0.38, 5.7, 0.36, item, 11, color=color)
    elif color == C_WHITE:
        pass
    else:
        bg_c = RGBColor(0xEE,0xF0,0xFF) if color == C_PRIMARY else (
               RGBColor(0xF0,0xFF,0xF4) if color == C_ACCENT else C_LIGHT)
        rect(sl, 7.0, 1.82+j*0.38, 5.7, 0.34, bg_c)
        tx(sl, 7.1, 1.83+j*0.38, 5.6, 0.32, item, 11, bold=True, color=color)
footer(sl)


# ════════════════════════════════════════════
# Slide 9: AI活用の役割分担と注意点
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "05  AI活用の役割分担と注意点")

rect(sl, 0.4, 1.25, 5.8, 3.6, C_WHITE)
rect(sl, 0.4, 1.25, 5.8, 0.5, C_DARK)
tx(sl, 0.4, 1.25, 5.8, 0.5, "  人間（SE）が担当", 13, bold=True, color=C_WHITE)
human_items = [
    "要件定義・課題の言語化",
    "UI/UXの方向性決定",
    "実装レビュー・動作確認",
    "インフラ構成の意思決定",
    "セキュリティ方針の決定",
    "dev/mainブランチ管理の判断",
]
for j, item in enumerate(human_items):
    tx(sl, 0.6, 1.85+j*0.4, 5.4, 0.38, f"✓  {item}", 12, color=C_TEXT)

rect(sl, 6.9, 1.25, 6.0, 3.6, C_WHITE)
rect(sl, 6.9, 1.25, 6.0, 0.5, C_PRIMARY)
tx(sl, 6.9, 1.25, 6.0, 0.5, "  AI（Claude）が担当", 13, bold=True, color=C_WHITE)
ai_items = [
    "コード実装（全量）",
    "DBスキーマ設計",
    "REST API 設計・実装",
    "テスト作成（77件）",
    "仕様書・ドキュメント作成",
    "Git コミット・ブランチ管理",
]
for j, item in enumerate(ai_items):
    tx(sl, 7.1, 1.85+j*0.4, 5.7, 0.38, f"✓  {item}", 12, color=C_TEXT)

# 重要な注意点
rect(sl, 0.4, 5.05, 12.5, 1.8, RGBColor(0xFF,0xF3,0xCD))
rect(sl, 0.4, 5.05, 0.18, 1.8, C_WARN)
tx(sl, 0.65, 5.1, 12.1, 0.45, "重要：「要件定義なき AI 活用」は失敗する", 14, bold=True, color=RGBColor(0x85,0x59,0x00))
tx(sl, 0.65, 5.58, 12.1, 0.4,
   "AI はあくまで「言われたことを実装するツール」。何を作るかが曖昧なまま指示すると、動くけど使えないシステムができる（GIGO原則）。",
   11, color=C_TEXT)
tx(sl, 0.65, 6.02, 12.1, 0.4,
   "今回の開発でも、要件定義・レビュー・意思決定はすべて人間が実施。この役割を果たせたのは SE としての知識と経験があったから。",
   11, color=C_TEXT)
footer(sl)


# ════════════════════════════════════════════
# Slide 10: 素人がAIだけで作る危険性
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, RGBColor(0xFF,0xF8,0xF8))
header(sl, "06  素人が AI だけで作る危険性", C_DANGER)

tx(sl, 0.5, 1.2, 12, 0.45,
   "AI は「動くコード」を生成できるが、「安全で正しいシステム」を保証しない。技術知識のないまま AI 任せにすると以下のリスクが生じる。",
   11, color=C_TEXT)

risks = [
    ("セキュリティ脆弱性",
     "SQL インジェクション / XSS / 認証バイパスなど OWASP Top 10 の脆弱性を含むコードが生成されることがある。\n指摘しなければそのまま本番に出る。",
     C_DANGER),
    ("機密情報の漏洩",
     ".env や API キーをコードにハードコードしてしまう。Git に push して公開リポジトリに上げると即アウト。\n今回も .gitignore の設定を明示的に指示した。",
     C_DANGER),
    ("設計の破綻",
     "「とりあえず動く」設計になりやすく、後から機能追加・修正が困難になる。\nDB設計・API設計の妥当性は技術者がレビューしなければわからない。",
     C_WARN),
    ("AI の幻覚（ハルシネーション）",
     "存在しないライブラリ・メソッドを自信満々に使用するコードを生成することがある。\n動かして初めて気づく。テストやレビューなしでは発見できない。",
     C_WARN),
    ("要件の曖昧さによる手戻り",
     "「いい感じに作って」という指示では意図と異なるものができる。\n要件を正確に言語化できなければ、何度作り直しても同じ結果になる。",
     RGBColor(0x85,0x59,0x00)),
    ("運用・保守の困難",
     "ログ・監視・デプロイ手順が整備されないまま本番稼働。障害発生時に対処できない。\n今回も systemd・Cloudflare Tunnel・自動デプロイを設計から考えた。",
     RGBColor(0x85,0x59,0x00)),
]
for j, (title, desc, color) in enumerate(risks):
    col = j % 2
    row = j // 2
    lx = 0.4 + col*6.45
    ty = 1.75 + row*1.55
    bg_c = RGBColor(0xFF,0xEE,0xEE) if color == C_DANGER else (
           RGBColor(0xFF,0xF8,0xEA) if color == C_WARN else RGBColor(0xFF,0xF8,0xEA))
    rect(sl, lx, ty, 6.1, 1.4, bg_c)
    rect(sl, lx, ty, 6.1, 0.4, color)
    tx(sl, lx+0.1, ty+0.05, 5.9, 0.33, title, 12, bold=True, color=C_WHITE)
    tx(sl, lx+0.1, ty+0.48, 5.9, 0.85, desc, 9, color=C_TEXT)

footer(sl, "素人のAI活用は「動くゴミ」を量産するリスクがある")


# ════════════════════════════════════════════
# Slide 11: 工数削減効果
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_LIGHT)
header(sl, "07  工数削減効果", C_ACCENT)

rect(sl, 0.4, 1.25, 5.9, 4.1, C_WHITE)
rect(sl, 0.4, 1.25, 5.9, 0.5, C_MUTED)
tx(sl, 0.4, 1.25, 5.9, 0.5, "  従来（個人開発のみ）", 13, bold=True, color=C_WHITE)
tx(sl, 1.5, 2.05, 3.2, 1.3, "3ヶ月以上", 40, bold=True, color=C_DANGER, align=PP_ALIGN.CENTER)
for j, line in enumerate(["設計: 2〜3週間", "実装: 6〜8週間", "テスト: 2週間", "ドキュメント: 1週間"]):
    tx(sl, 0.6, 3.5+j*0.35, 5.5, 0.32, f"• {line}", 11, color=C_TEXT)

tx(sl, 6.2, 3.6, 0.9, 0.6, "=>", 22, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

rect(sl, 7.0, 1.25, 6.0, 4.1, C_WHITE)
rect(sl, 7.0, 1.25, 6.0, 0.5, C_ACCENT)
tx(sl, 7.0, 1.25, 6.0, 0.5, "  AI活用（Claude × SE）", 13, bold=True, color=C_WHITE)
tx(sl, 8.2, 2.05, 3.2, 1.3, "約  3日", 40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
for j, line in enumerate(["実装: AIが全量担当", "テスト: 77件自動生成", "ドキュメント: 自動生成", "コードレビュー: 即時"]):
    tx(sl, 7.2, 3.5+j*0.35, 5.7, 0.32, f"• {line}", 11, color=C_TEXT)

rect(sl, 0.4, 5.55, 12.5, 1.25, C_PRIMARY)
tx(sl, 0.4, 5.58, 12.5, 0.52, "工数削減率：約 97 ％  （3ヶ月以上 → 約3日）", 26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
tx(sl, 0.4, 6.1, 12.5, 0.45,
   "人間の作業は「何を作るか」の定義とレビューに集中  ／  実装・テスト・ドキュメントは AI が担当",
   12, color=RGBColor(0xCC,0xCF,0xFF), align=PP_ALIGN.CENTER)
footer(sl)


# ════════════════════════════════════════════
# Slide 12: まとめ・今後の展望
# ════════════════════════════════════════════
sl = add_slide()
bg(sl, C_DARK)
header(sl, "08  まとめ・今後の展望")

rect(sl, 0.4, 1.25, 5.9, 5.6, RGBColor(0x36,0x39,0x3F))
tx(sl, 0.5, 1.3, 5.7, 0.45, "今回の成果", 13, bold=True, color=C_ACCENT)
summary_items = [
    "旧JSON管理 → MySQL+API+Dashboard",
    "への全面刷新を約3日で完了",
    "",
    "7項目の改修要件をすべて実装",
    "",
    "APIテスト 77件 をパス",
    "",
    "チャンネル別ローテーション管理を",
    "odai_usageテーブルで実現",
    "",
    "SE知識があったからこそ成立した開発",
]
for j, line in enumerate(summary_items):
    tx(sl, 0.6, 1.85+j*0.38, 5.5, 0.36, line, 11, color=C_WHITE)

rect(sl, 6.9, 1.25, 6.0, 5.6, RGBColor(0x36,0x39,0x3F))
tx(sl, 7.0, 1.3, 5.8, 0.45, "今後の展望", 13, bold=True, color=C_PRIMARY)
future_items = [
    "本番環境へのデプロイ",
    "（Cloudflare Tunnel HTTPS化）",
    "",
    "GitHub Webhook 自動デプロイ整備",
    "",
    "dev → main ブランチマージ",
    "",
    "複数サーバーへの展開",
    "",
    "業務への AI 活用拡大を検討",
    "（チームでの活用事例を模索）",
]
for j, line in enumerate(future_items):
    tx(sl, 7.1, 1.85+j*0.38, 5.7, 0.36, line, 11, color=C_WHITE)

rect(sl, 0.4, 7.05, 12.5, 0.35, C_PRIMARY)
tx(sl, 0.4, 7.07, 12.5, 0.3,
   "「何を作るか」を定義できる人間  ×  「速く作る」AI  =  最大の生産性",
   12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
output = "report.pptx"
prs.save(output)
print(f"Saved: {output}")
